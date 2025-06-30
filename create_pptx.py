import os
import re
from pptx import Presentation
from pptx.util import Inches

# 创建一个 PowerPoint 演示文稿
presentation = Presentation()

# 设置幻灯片的页面大小为 16:9 比例（宽度 13.33 英寸，高度 7.5 英寸）
presentation.slide_width = Inches(13.33)
presentation.slide_height = Inches(7.5)

# 目标文件夹路径
image_folder = './result'

# 提取文件名中的数字并排序
def extract_number(file_name):
    # 使用正则表达式提取文件名中的数字部分
    match = re.search(r'(\d+)', file_name)
    return int(match.group(0)) if match else float('inf')  # 如果没有数字，返回一个很大的数

# 获取所有图片文件，并按数字顺序排序
image_files = [f for f in os.listdir(image_folder) if f.endswith(('.png', '.jpg'))]
image_files.sort(key=extract_number)  # 按数字排序

# 遍历排序后的图片文件
for file_name in image_files:
    try:
        # 创建一个新幻灯片
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # 选择空白布局

        # 获取图片路径
        image_path = os.path.join(image_folder, file_name)

        # 添加图片到幻灯片
        slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=Inches(13.33), height=Inches(7.5))

        print(f"图片 {file_name} 已添加到幻灯片.")

    except Exception as e:
        print(f"处理图片 {file_name} 时发生错误: {e}")

# 保存 PowerPoint 文件
pptx_file_path = './result/combined_presentation.pptx'
presentation.save(pptx_file_path)
print(f"PPT 文件已保存为: {pptx_file_path}")
